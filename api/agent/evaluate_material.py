import os
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import nltk
from nltk.tokenize import sent_tokenize
from langchain.vectorstores import FAISS
from langchain.embeddings import HuggingFaceEmbeddings
import api.models.vector_store as vs

nltk.download("punkt")

vector_db = vs.load_proposal_vector_db()

# === [1] 문서 텍스트 추출 함수 ===
def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    elif ext == ".pdf":
        text = ""
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        return text
    elif ext == ".docx":
        doc = DocxDocument(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    elif ext == ".pptx":
        prs = Presentation(file_path)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    else:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {ext}")

# === [2] 문장 유사도 분석 ===
def analyze_similarity(file_path: str, vector_db_path: str, top_k: int = 3):
    text = extract_text(file_path)
    sentences = sent_tokenize(text)

    # 벡터 DB 로딩
    embedding_model = HuggingFaceEmbeddings(model_name="jhgan/ko-sbert-nli")

    scored_sentences = []
    for idx, sentence in enumerate(sentences):
        results = vector_db.similarity_search_with_score(sentence, k=top_k)
        if results:
            top_score = results[0][1]  # 두 번째 요소가 점수
            scored_sentences.append((idx, sentence, top_score))

    # 유사도 평균
    avg_score = sum([score for _, _, score in scored_sentences]) / len(scored_sentences)

    # 상위/하위 문장 (유사도 기준)
    scored_sentences.sort(key=lambda x: x[2])
    lowest = scored_sentences[:3]
    highest = scored_sentences[-3:]

    return {
        "평균 유사도": avg_score,
        "유사도 높은 문장": [{"문장번호": i, "문장": s, "점수": sc} for i, s, sc in reversed(highest)],
        "유사도 낮은 문장": [{"문장번호": i, "문장": s, "점수": sc} for i, s, sc in lowest]
    }
